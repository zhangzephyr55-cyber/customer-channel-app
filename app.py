import sqlite3, os, io, json
from datetime import datetime, date, timedelta
from flask import Flask, render_template, request, jsonify, send_file
import openpyxl

app = Flask(__name__)
app.secret_key = os.urandom(24).hex()
ADMIN_PASSWORD = 'admin123'  # 首次使用请修改此密码

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
DB_PATH = os.path.join(BASE_DIR, 'data.db')
UPLOAD_FOLDER = os.path.join(BASE_DIR, 'uploads')
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

CURRENT_DATE = date(2026, 7, 8)
ONE_YEAR_AGO = CURRENT_DATE - timedelta(days=365)  # 2025-07-08

CHANNEL_NAMES = ['商超终端', '批发', 'CS', '电商', '网点', '零食折扣', '即时零售', '其他']
CHANNEL_KEYS = ['supermarket', 'wholesale', 'cs', 'ecommerce', 'outlet', 'snack_discount', 'instant_retail', 'other']

def get_db():
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    conn.execute("PRAGMA journal_mode=WAL")
    return conn

def init_db():
    conn = get_db()
    conn.executescript("""
        CREATE TABLE IF NOT EXISTS transactions (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            date TEXT, doc_no TEXT, customer_code TEXT,
            customer_name TEXT, salesperson TEXT,
            region TEXT, province TEXT, amount REAL
        );
        CREATE TABLE IF NOT EXISTS channel_fillings (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            customer_code TEXT, salesperson TEXT,
            supermarket REAL DEFAULT 0, wholesale REAL DEFAULT 0,
            cs REAL DEFAULT 0, ecommerce REAL DEFAULT 0,
            outlet REAL DEFAULT 0, snack_discount REAL DEFAULT 0,
            instant_retail REAL DEFAULT 0, other REAL DEFAULT 0,
            updated_at TEXT DEFAULT (datetime('now','localtime')),
            UNIQUE(customer_code, salesperson)
        );
    """)
    conn.commit()
    conn.close()

init_db()

def parse_date(val):
    if isinstance(val, (datetime, date)):
        return val.strftime('%Y-%m-%d')
    if isinstance(val, (int, float)):
        base = datetime(1899, 12, 30)
        return (base + timedelta(days=int(val))).strftime('%Y-%m-%d')
    s = str(val).strip()
    for fmt in ['%Y-%m-%d', '%Y/%m/%d', '%Y.%m.%d', '%Y年%m月%d日']:
        try:
            return datetime.strptime(s, fmt).strftime('%Y-%m-%d')
        except:
            pass
    return s

def get_effective_customers(salesperson=None):
    """SQL param order: [?1:year_ago, ?2:last_start, ?3:last_end, ?4:year_start, ?5:salesperson]"""
    conn = get_db()
    yr_start = str(CURRENT_DATE.year) + '-01-01'
    last_start = str(CURRENT_DATE.year - 1) + '-01-01'
    last_end = str(CURRENT_DATE.year - 1) + '-12-31'
    year_ago = str(ONE_YEAR_AGO)

    if salesperson:
        rows = conn.execute("""
            SELECT t.customer_code, t.customer_name, t.salesperson, t.region, t.province,
                   COUNT(DISTINCT CASE WHEN t.date >= ? THEN t.doc_no ELSE NULL END) AS order_count_last_year,
                   COALESCE(SUM(CASE WHEN t.date >= ? AND t.date <= ? THEN t.amount ELSE 0 END),0) AS last_year_sales,
                   COALESCE(SUM(CASE WHEN t.date >= ? THEN t.amount ELSE 0 END),0) AS this_year_sales
            FROM transactions t
            WHERE t.salesperson = ?
            GROUP BY t.customer_code
            HAVING order_count_last_year >= 3
            ORDER BY t.customer_name
        """, [year_ago, last_start, last_end, yr_start, salesperson]).fetchall()
    else:
        rows = conn.execute("""
            SELECT t.customer_code, t.customer_name, t.salesperson, t.region, t.province,
                   COUNT(DISTINCT CASE WHEN t.date >= ? THEN t.doc_no ELSE NULL END) AS order_count_last_year,
                   COALESCE(SUM(CASE WHEN t.date >= ? AND t.date <= ? THEN t.amount ELSE 0 END),0) AS last_year_sales,
                   COALESCE(SUM(CASE WHEN t.date >= ? THEN t.amount ELSE 0 END),0) AS this_year_sales
            FROM transactions t
            GROUP BY t.customer_code
            HAVING order_count_last_year >= 3
            ORDER BY t.customer_name
        """, [year_ago, last_start, last_end, yr_start]).fetchall()

    result = []
    for row in rows:
        d = dict(row)
        if salesperson:
            f = conn.execute("SELECT * FROM channel_fillings WHERE customer_code=? AND salesperson=?", [d['customer_code'], salesperson]).fetchone()
        else:
            f = conn.execute("SELECT * FROM channel_fillings WHERE customer_code=? ORDER BY updated_at DESC LIMIT 1", [d['customer_code']]).fetchone()
        if f:
            d.update({k: float(f[k]) for k in CHANNEL_KEYS})
        else:
            d.update({k: 0.0 for k in CHANNEL_KEYS})
        result.append(d)
    conn.close()
    return result

@app.route('/')
def index():
    return render_template('index.html', year=CURRENT_DATE.year, channels=list(zip(CHANNEL_NAMES, CHANNEL_KEYS)))

@app.route('/api/customers')
def api_customers():
    name = request.args.get('name', '').strip()
    if not name:
        return jsonify({'error': '请输入姓名'}), 400
    customers = get_effective_customers(name)
    return jsonify({'customers': customers, 'salesperson': name})

@app.route('/api/save', methods=['POST'])
def api_save():
    data = request.get_json()
    if not data:
        return jsonify({'error': '无效数据'}), 400
    sp = data.get('salesperson', '').strip()
    channels = data.get('channels', [])
    if not sp or not channels:
        return jsonify({'error': '缺少必要参数'}), 400
    conn = get_db()
    for ch in channels:
        conn.execute("""
            INSERT INTO channel_fillings (customer_code, salesperson, supermarket, wholesale, cs, ecommerce, outlet, snack_discount, instant_retail, other, updated_at)
            VALUES (?,?,?,?,?,?,?,?,?,?, datetime('now','localtime'))
            ON CONFLICT(customer_code, salesperson) DO UPDATE SET
                supermarket=excluded.supermarket, wholesale=excluded.wholesale,
                cs=excluded.cs, ecommerce=excluded.ecommerce,
                outlet=excluded.outlet, snack_discount=excluded.snack_discount,
                instant_retail=excluded.instant_retail, other=excluded.other,
                updated_at=datetime('now','localtime')
        """, [ch['customer_code'], sp,
              float(ch.get('supermarket',0)), float(ch.get('wholesale',0)),
              float(ch.get('cs',0)), float(ch.get('ecommerce',0)),
              float(ch.get('outlet',0)), float(ch.get('snack_discount',0)),
              float(ch.get('instant_retail',0)), float(ch.get('other',0))])
    conn.commit()
    conn.close()
    return jsonify({'success': True})


@app.route('/admin/login', methods=['POST'])
def admin_login():
    pwd = request.form.get('password', '')
    if pwd == ADMIN_PASSWORD:
        from flask import session
        session['admin_auth'] = True
        return jsonify({'success': True})
    return jsonify({'error': '密码错误'}), 403

@app.route('/admin/logout')
def admin_logout():
    from flask import session
    session.pop('admin_auth', None)
    return ('', 204)

@app.route('/admin')
def admin():
    from flask import session, redirect
    if not session.get('admin_auth'):
        return render_template('admin_login.html')
    return render_template('admin.html', year=CURRENT_DATE.year)

@app.route('/admin/upload', methods=['POST'])
def admin_upload():
    if 'file' not in request.files:
        return jsonify({'error': '请选择文件'}), 400
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': '请选择文件'}), 400
    try:
        wb = openpyxl.load_workbook(file, data_only=True)
        ws = wb.active
        headers = [cell.value for cell in ws[1]]
        expected = ['年月日','单据号','客户代码','客户','业务员','大区','省份','销售金额']
        hmap = {}
        for h in expected:
            for i, ch in enumerate(headers):
                if ch and str(ch).strip() == h:
                    hmap[h] = i
                    break
        missing = [h for h in expected if h not in hmap]
        if missing:
            return jsonify({'error': f'Excel缺少列: {", ".join(missing)}'}), 400
        txns = []
        for row in ws.iter_rows(min_row=2, values_only=True):
            if row[hmap['客户代码']] is None:
                continue
            txns.append({
                'date': parse_date(row[hmap['年月日']]),
                'doc_no': str(row[hmap['单据号']] or '').strip(),
                'customer_code': str(row[hmap['客户代码']] or '').strip(),
                'customer_name': str(row[hmap['客户']] or '').strip(),
                'salesperson': str(row[hmap['业务员']] or '').strip(),
                'region': str(row[hmap['大区']] or '').strip(),
                'province': str(row[hmap['省份']] or '').strip(),
                'amount': float(row[hmap['销售金额']] or 0)
            })
        if not txns:
            return jsonify({'error': '未读取到有效数据'}), 400
        conn = get_db()
        conn.execute('DELETE FROM transactions')
        conn.execute('DELETE FROM channel_fillings')
        conn.executemany("""
            INSERT INTO transactions (date, doc_no, customer_code, customer_name, salesperson, region, province, amount)
            VALUES (:date, :doc_no, :customer_code, :customer_name, :salesperson, :region, :province, :amount)
        """, txns)
        conn.commit()
        eff = conn.execute("""
            SELECT COUNT(*) FROM (
                SELECT customer_code FROM transactions
                GROUP BY customer_code
                HAVING COUNT(DISTINCT CASE WHEN date>=? THEN doc_no ELSE NULL END) >= 3
            )
        """, [str(ONE_YEAR_AGO)]).fetchone()[0]
        sp_count = conn.execute("SELECT COUNT(DISTINCT salesperson) FROM transactions").fetchone()[0]
        filled = conn.execute("""SELECT COUNT(DISTINCT customer_code) FROM channel_fillings WHERE ROUND(supermarket+wholesale+cs+ecommerce+outlet+snack_discount+instant_retail+other,1) = 100""").fetchone()[0]
        conn.close()
        return jsonify({
            'success': True,
            'message': f'上传成功！共导入 {len(txns)} 条记录',
            'stats': {'total': len(txns), 'effective': eff, 'salespeople': sp_count, 'filled': filled}
        })
    except Exception as e:
        return jsonify({'error': f'处理文件出错: {str(e)}'}), 500


@app.route('/admin/import', methods=['POST'])
def admin_import():
    from flask import session
    if not session.get('admin_auth'):
        return jsonify({'error': '未登录'}), 403
    if 'file' not in request.files:
        return jsonify({'error': '请选择文件'}), 400
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': '请选择文件'}), 400
    try:
        wb = openpyxl.load_workbook(file, data_only=True)
        if '渠道备份' not in wb.sheetnames:
            return jsonify({'error': '此文件不是有效的备份文件（缺少 渠道备份 工作表）'}), 400
        ws = wb['渠道备份']
        headers = [cell.value for cell in ws[1]]
        expected = ['客户代码', '业务员'] + CHANNEL_KEYS
        hmap = {}
        for h in expected:
            for i, ch in enumerate(headers):
                if ch and str(ch).strip() == h:
                    hmap[h] = i; break
        missing = [h for h in expected if h not in hmap]
        if missing:
            return jsonify({'error': f'备份文件格式不匹配，缺少: {", ".join(missing)}'}), 400
        conn = get_db()
        restored = 0
        for row in ws.iter_rows(min_row=2, values_only=True):
            cc = str(row[hmap['客户代码']] or '').strip()
            sp = str(row[hmap['业务员']] or '').strip()
            if not cc or not sp:
                continue
            vals = {}
            for k in CHANNEL_KEYS:
                vals[k] = float(row[hmap[k]] or 0)
            conn.execute("""
                INSERT INTO channel_fillings (customer_code, salesperson, supermarket, wholesale, cs, ecommerce, outlet, snack_discount, instant_retail, other, updated_at)
                VALUES (?,?,?,?,?,?,?,?,?,?, datetime('now','localtime'))
                ON CONFLICT(customer_code, salesperson) DO UPDATE SET
                    supermarket=excluded.supermarket, wholesale=excluded.wholesale,
                    cs=excluded.cs, ecommerce=excluded.ecommerce,
                    outlet=excluded.outlet, snack_discount=excluded.snack_discount,
                    instant_retail=excluded.instant_retail, other=excluded.other,
                    updated_at=datetime('now','localtime')
            """, [cc, sp, vals['supermarket'], vals['wholesale'], vals['cs'], vals['ecommerce'],
                  vals['outlet'], vals['snack_discount'], vals['instant_retail'], vals['other']])
            restored += 1
        conn.commit()
        conn.close()
        return jsonify({'success': True, 'message': f'恢复成功！共恢复 {restored} 家客户的渠道数据'})
    except Exception as e:
        return jsonify({'error': f'恢复失败: {str(e)}'}), 500

@app.route('/admin/export')
def admin_export():
    customers = get_effective_customers()
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = '客户渠道汇总'
    headers = ['客户代码','客户','业务员','大区','省份','上一年度销售收入','本年度销售收入']
    headers += [f'{n}占比(%)' for n in CHANNEL_NAMES]
    headers.append('是否有电商')
    for ci, h in enumerate(headers, 1):
        c = ws.cell(row=1, column=ci, value=h)
        c.font = openpyxl.styles.Font(bold=True, size=11)
        c.alignment = openpyxl.styles.Alignment(horizontal='center')
        c.border = openpyxl.styles.Border(bottom=openpyxl.styles.Side(style='thin'))
    for ri, cust in enumerate(customers, 2):
        ws.cell(row=ri, column=1, value=cust['customer_code'])
        ws.cell(row=ri, column=2, value=cust['customer_name'])
        ws.cell(row=ri, column=3, value=cust['salesperson'])
        ws.cell(row=ri, column=4, value=cust['region'])
        ws.cell(row=ri, column=5, value=cust['province'])
        ws.cell(row=ri, column=6, value=round(cust.get('last_year_sales',0),2))
        ws.cell(row=ri, column=7, value=round(cust.get('this_year_sales',0),2))
        has_ec = False
        for ji, k in enumerate(CHANNEL_KEYS):
            v = round(float(cust.get(k,0)), 1)
            cell = ws.cell(row=ri, column=8+ji, value=v)
            cell.number_format = '0.0'
            if k == 'ecommerce' and v > 0:
                has_ec = True
        ws.cell(row=ri, column=8+len(CHANNEL_NAMES), value='是' if has_ec else '否')
    for ci in range(1, len(headers)+1):
        mx = max(len(str(ws.cell(row=r,column=ci).value or '')) for r in range(1, len(customers)+2))
        ws.column_dimensions[openpyxl.utils.get_column_letter(ci)].width = max(mx+3, 12)
        # Add backup sheet with channel fillings data
    ws2 = wb.create_sheet('渠道备份')
    bk_headers = ['客户代码','客户','业务员','大区','省份'] + CHANNEL_KEYS
    for ci, h in enumerate(bk_headers, 1):
        c = ws2.cell(row=1, column=ci, value=h)
        c.font = openpyxl.styles.Font(bold=True, size=11)
    # Get all channel fillings with customer info
    conn2 = get_db()
    fill_rows = conn2.execute("""
        SELECT f.customer_code, t.customer_name, f.salesperson, t.region, t.province,
               f.supermarket, f.wholesale, f.cs, f.ecommerce, f.outlet, f.snack_discount, f.instant_retail, f.other
        FROM channel_fillings f
        LEFT JOIN (SELECT DISTINCT customer_code, customer_name, region, province FROM transactions) t
        ON f.customer_code = t.customer_code
        ORDER BY f.salesperson, f.customer_code
    """).fetchall()
    conn2.close()
    for ri, row in enumerate(fill_rows, 2):
        ws2.cell(row=ri, column=1, value=row['customer_code'])
        ws2.cell(row=ri, column=2, value=row['customer_name'])
        ws2.cell(row=ri, column=3, value=row['salesperson'])
        ws2.cell(row=ri, column=4, value=row['region'])
        ws2.cell(row=ri, column=5, value=row['province'])
        for ji, k in enumerate(CHANNEL_KEYS):
            ws2.cell(row=ri, column=6+ji, value=round(float(row[k]), 1))
    out = io.BytesIO()
    wb.save(out)
    out.seek(0)
    return send_file(out, mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                     as_attachment=True, download_name=f'客户渠道汇总_{CURRENT_DATE.year}.xlsx')

@app.route('/admin/stats')
def admin_stats():
    conn = get_db()
    eff = conn.execute("""
        SELECT COUNT(*) FROM (
            SELECT customer_code FROM transactions
            GROUP BY customer_code
            HAVING COUNT(DISTINCT CASE WHEN date>=? THEN doc_no ELSE NULL END) >= 3
        )
    """, [str(ONE_YEAR_AGO)]).fetchone()[0]
    filled = conn.execute("""SELECT COUNT(DISTINCT customer_code) FROM channel_fillings WHERE ROUND(supermarket+wholesale+cs+ecommerce+outlet+snack_discount+instant_retail+other,1) = 100""").fetchone()[0]
    sp_rows = conn.execute("SELECT DISTINCT salesperson FROM transactions ORDER BY salesperson").fetchall()
    pstats = []
    for sp in sp_rows:
        sp_name = sp['salesperson']
        total = conn.execute("""
            SELECT COUNT(*) FROM (
                SELECT customer_code FROM transactions WHERE salesperson=?
                GROUP BY customer_code
                HAVING COUNT(DISTINCT CASE WHEN date>=? THEN doc_no ELSE NULL END) >= 3
            )
        """, [sp_name, str(ONE_YEAR_AGO)]).fetchone()[0]
        sp_filled = conn.execute("""
            SELECT COUNT(DISTINCT f.customer_code) FROM channel_fillings f
            WHERE f.salesperson=? AND ROUND(f.supermarket+f.wholesale+f.cs+f.ecommerce+f.outlet+f.snack_discount+f.instant_retail+f.other,1) = 100
        """, [sp_name]).fetchone()[0]
        pstats.append({'name': sp_name, 'total': total, 'filled': sp_filled, 'remaining': total-sp_filled})
    conn.close()
    return jsonify({'effective': eff, 'filled': filled, 'remaining': eff-filled, 'person_stats': pstats})


@app.route('/admin/analysis/data')
def admin_analysis_data():
    conn = get_db()
    customers = get_effective_customers()
    total = len(customers)
    filled_customers = [c for c in customers if abs(sum(float(c.get(k,0)) for k in CHANNEL_KEYS) - 100) < 0.1]
    filled = len(filled_customers)
    
    # Channel totals & averages
    ch_data = {}
    for k, n in zip(CHANNEL_KEYS, CHANNEL_NAMES):
        vals = [float(c.get(k,0)) for c in filled_customers]
        ch_data[n] = {
            'total': round(sum(vals), 1),
            'avg': round(sum(vals)/len(vals), 1) if vals else 0,
            'count_gt0': sum(1 for v in vals if v > 0),
            'pct_of_customers': round(sum(1 for v in vals if v > 0)/total*100, 1) if total else 0
        }
    
    # Dominant channel (which channel has highest % per customer)
    dominant = {}
    for c in filled_customers:
        max_val = max((float(c.get(k,0)), k) for k in CHANNEL_KEYS)
        if max_val[0] > 0:
            n = CHANNEL_NAMES[CHANNEL_KEYS.index(max_val[1])]
            dominant[n] = dominant.get(n, 0) + 1
    
    # Region breakdown
    regions = {}
    for c in customers:
        r = c.get('region', '未知') or '未知'
        if r not in regions:
            regions[r] = {'total': 0, 'filled': 0}
        regions[r]['total'] += 1
        if abs(sum(float(c.get(k,0)) for k in CHANNEL_KEYS) - 100) < 0.1:
            regions[r]['filled'] += 1
    
    # Ecommerce summary
    has_ec = sum(1 for c in filled_customers if float(c.get('ecommerce',0)) > 0)
    
    conn.close()
    return jsonify({
        'total_customers': total,
        'filled_customers': filled,
        'completion_pct': round(filled/total*100, 1) if total else 0,
        'channels': ch_data,
        'dominant_channel': dominant,
        'regions': regions,
        'has_ecommerce': has_ec
    })

@app.route('/admin/analysis/export-ppt')
def admin_analysis_export_ppt():
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    import math
    
    conn = get_db()
    customers = get_effective_customers()
    filled_customers = [c for c in customers if abs(sum(float(c.get(k,0)) for k in CHANNEL_KEYS) - 100) < 0.1]
    conn.close()
    
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg = sl.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(0, 47, 167)
    txBox = sl.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = '客户渠道数据分析报告'
    p.font.size = Pt(36)
    p.font.color.rgb = RGBColor(255,255,255)
    p.font.bold = True
    p2 = tf.add_paragraph()
    p2.text = f'报告日期：{CURRENT_DATE.year}年{CURRENT_DATE.month}月{CURRENT_DATE.day}日'
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(200,200,255)
    p3 = tf.add_paragraph()
    p3.text = f'有效客户：{len(customers)} 家 | 已填写：{len(filled_customers)} 家'
    p3.font.size = Pt(16)
    p3.font.color.rgb = RGBColor(200,200,255)
    
    # Slide 2: Channel distribution table
    sl2 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox2 = sl2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf2 = txBox2.text_frame
    tf2.paragraphs[0].text = '渠道分布概览'
    tf2.paragraphs[0].font.size = Pt(24)
    tf2.paragraphs[0].font.bold = True
    
    # Table
    rows_n = 10
    cols_n = 5
    table = sl2.shapes.add_table(rows_n, cols_n, Inches(0.5), Inches(1.2), Inches(10), Inches(5)).table
    headers_t = ['渠道', '总占比', '平均占比', '覆盖客户数', '覆盖率(%)']
    for ci, h in enumerate(headers_t):
        cell = table.cell(0, ci)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255,255,255)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0,47,167)
    
    for ri, (k, n) in enumerate(zip(CHANNEL_KEYS, CHANNEL_NAMES)):
        vals = [float(c.get(k,0)) for c in filled_customers]
        avg = round(sum(vals)/len(vals), 1) if vals else 0
        total = round(sum(vals), 1)
        cnt = sum(1 for v in vals if v > 0)
        pct = round(cnt/len(customers)*100, 1) if customers else 0
        table.cell(ri+1, 0).text = n
        table.cell(ri+1, 1).text = str(total)
        table.cell(ri+1, 2).text = str(avg)
        table.cell(ri+1, 3).text = str(cnt)
        table.cell(ri+1, 4).text = str(pct)
        for ci in range(5):
            cell = table.cell(ri+1, ci)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                if ci >= 1:
                    p.alignment = 3  # right
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(240, 242, 250)
    
    # Summary row
    table.cell(9, 0).text = '合计'
    for p in table.cell(9, 0).text_frame.paragraphs:
        p.font.bold = True
    table.cell(9, 1).text = str(round(sum(float(c.get(k,0)) for c in filled_customers for k in CHANNEL_KEYS), 1))
    for p in table.cell(9, 1).text_frame.paragraphs:
        p.font.bold = True
    
    # Slide 3: Summary insights
    sl3 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox3 = sl3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf3 = txBox3.text_frame
    tf3.paragraphs[0].text = '关键洞察'
    tf3.paragraphs[0].font.size = Pt(24)
    tf3.paragraphs[0].font.bold = True
    
    insights = []
    insights.append(f'有效客户总数：{len(customers)} 家')
    insights.append(f'已填写渠道数据：{len(filled_customers)} 家（{round(len(filled_customers)/len(customers)*100,1) if customers else 0}%）')
    ec_count = sum(1 for c in filled_customers if float(c.get('ecommerce',0)) > 0)
    insights.append(f'有电商业务的客户：{ec_count} 家')
    
    top_ch = ''
    max_dom = 0
    dominant = {}
    for c in filled_customers:
        max_v = max((float(c.get(k,0)), k) for k in CHANNEL_KEYS)
        if max_v[0] > 0:
            n = CHANNEL_NAMES[CHANNEL_KEYS.index(max_v[1])]
            dominant[n] = dominant.get(n, 0) + 1
            if dominant[n] > max_dom:
                max_dom = dominant[n]
                top_ch = n
    if top_ch:
        insights.append(f'主力渠道：{top_ch}（{max_dom} 家客户以此为主渠道）')
    
    txBody = sl3.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(10), Inches(5))
    tfBody = txBody.text_frame
    tfBody.word_wrap = True
    for i, ins in enumerate(insights):
        if i == 0:
            p = tfBody.paragraphs[0]
        else:
            p = tfBody.add_paragraph()
        p.text = f'• {ins}'
        p.font.size = Pt(16)
        p.space_after = Pt(12)
    
    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return send_file(out, mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                     as_attachment=True, download_name=f'客户渠道分析报告_{CURRENT_DATE.year}.pptx')

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=int(os.environ.get('PORT', 5000)), debug=False)





